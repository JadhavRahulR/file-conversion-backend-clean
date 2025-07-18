from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import sys
import os

def convert_pdf_to_pptx(pdf_path, output_path):
    # Step 1: Convert PDF pages to images
    images = convert_from_path(pdf_path)

    # Step 2: Create a new PowerPoint presentation
    prs = Presentation()

    # Set slide size to 16:9 widescreen
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Step 3: Add each image to a new slide
    for index, image in enumerate(images):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        temp_image_path = f"temp_slide_{index}.png"
        image.save(temp_image_path, "PNG")

        slide.shapes.add_picture(
            temp_image_path,
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

        os.remove(temp_image_path)

    # Step 4: Save presentation
    prs.save(output_path)
    print(f"PPTX created: {output_path}")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python pdf_to_pptx.py input.pdf output.pptx")
        sys.exit(1)

    input_pdf = sys.argv[1]
    output_pptx = sys.argv[2]
    convert_pdf_to_pptx(input_pdf, output_pptx)
